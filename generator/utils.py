import os
import json
import copy
import tempfile
from pptx import Presentation
from pptx.text.text import _Paragraph
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

# --- HELPER 1: Get All Text ---
def get_all_text(shape):
    text = ""
    if shape.has_text_frame:
        text += "".join([p.text for p in shape.text_frame.paragraphs])
    if hasattr(shape, "shapes"): 
        for s in shape.shapes: text += get_all_text(s)
    return text

# --- HELPER 2: Text Replace ---
def replace_text(shape, tag, replacement):
    if hasattr(shape, "shapes"):
        for s in shape.shapes: replace_text(s, tag, replacement)
    if not shape.has_text_frame: return
    for p in shape.text_frame.paragraphs:
        if tag in p.text:
            replaced = False
            for run in p.runs:
                if tag in run.text:
                    run.text = run.text.replace(tag, replacement)
                    replaced = True
            if not replaced:
                p.text = p.text.replace(tag, replacement)

# --- HELPER 3: Bullets Replace ---
def replace_with_bullets(shape, tag, bullet_list):
    if hasattr(shape, "shapes"):
        for s in shape.shapes: replace_with_bullets(s, tag, bullet_list)
    if not shape.has_text_frame: return
    has_tag = False
    for p in shape.text_frame.paragraphs:
        if tag in p.text:
            has_tag = True
            break
    if not has_tag: return
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP 
    if not bullet_list:
        tf.clear()
        return
    template_p_xml = copy.deepcopy(tf.paragraphs[0]._p)
    tf.clear()
    p0 = tf.paragraphs[0]
    p0._p.getparent().replace(p0._p, template_p_xml)
    p0 = _Paragraph(template_p_xml, tf)
    if len(p0.runs) > 0:
        p0.runs[0].text = bullet_list[0]
        for i in range(len(p0.runs)-1, 0, -1): p0._p.remove(p0.runs[i]._r)
    else:
        p0.text = bullet_list[0]
    p0.space_before, p0.space_after = Pt(0), Pt(0)
    parent = p0._p.getparent()
    for point in bullet_list[1:]:
        new_p_xml = copy.deepcopy(template_p_xml)
        parent.append(new_p_xml)
        new_p = _Paragraph(new_p_xml, tf)
        if len(new_p.runs) > 0:
            new_p.runs[0].text = point
            for i in range(len(new_p.runs)-1, 0, -1): new_p._p.remove(new_p.runs[i]._r)
        else:
            new_p.text = point
        new_p.space_before, new_p.space_after = Pt(0), Pt(0)

# --- HELPER 4: Image Replacement ---
def replace_image(slide, shape, img_file):
    x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
    with tempfile.NamedTemporaryFile(delete=False) as t:
        for chunk in img_file.chunks(): t.write(chunk)
        p = t.name
    slide.shapes.add_picture(p, x, y, cx, cy)
    shape._element.getparent().remove(shape._element)
    os.unlink(p)

# --- HELPER 5: Native Slide Cloner ---
def clone_slide_to_presentation(source_prs, slide_index, target_prs):
    source_slide = source_prs.slides[slide_index]
    try:
        blank_layout = target_prs.slide_layouts[6] 
    except:
        blank_layout = target_prs.slide_layouts[0]
        
    new_slide = target_prs.slides.add_slide(blank_layout)
    
    if source_slide.background and source_slide.background.fill:
        try: new_slide.background.fill.copy_from(source_slide.background.fill)
        except: pass

    for shape in source_slide.shapes:
        if hasattr(shape, "shape_type") and shape.shape_type == 13:
            try:
                image_blob = shape.image.blob
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tf:
                    tf.write(image_blob)
                    tmp_path = tf.name
                new_pic = new_slide.shapes.add_picture(tmp_path, shape.left, shape.top, shape.width, shape.height)
                alt_text = shape._element._nvXxPr.cNvPr.attrib.get('descr', '')
                if alt_text: new_pic._element._nvXxPr.cNvPr.set('descr', alt_text)
                new_pic._element._nvXxPr.cNvPr.set('name', shape.name)
                os.unlink(tmp_path)
            except:
                new_el = copy.deepcopy(shape._element)
                new_slide.shapes._spTree.append(new_el)
        else:
            new_el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(new_el)
    return new_slide

# --- HELPER 6: HEX TO RGB ---
def hex_to_rgb(value):
    value = value.lstrip('#')
    return tuple(int(value[i:i+2], 16) for i in (0, 2, 4))

# --- HELPER 7: THE ULTIMATE XML COLOR HACK ---
# --- HELPER 7: THE ULTIMATE XML COLOR HACK (Updated for Shape 2 & Shape 4) ---
def apply_theme_color(shape, hex_color, force_child=False):
    try:
        is_target = force_child
        
        # 1. Target dhundo (Naam se: Shape 2, Shape 4, ya theme_accent)
        if hasattr(shape, "name"):
            s_name = shape.name.lower()
            # Yahan humne naye naam add kar diye hain
            if "theme_accent" in s_name or "shape 2" in s_name or "shape 4" in s_name:
                is_target = True
                
        try:
            if hasattr(shape, "_element") and hasattr(shape._element, "nvSpPr"):
                desc = shape._element.nvSpPr.cNvPr.attrib.get('descr', '').lower()
                if "theme_accent" in desc or "shape 2" in desc or "shape 4" in desc:
                    is_target = True
        except: pass

        # 2. RAW XML SURGERY (Yahan asli hack hai)
        if is_target:
            clean_hex = hex_color.lstrip('#').upper()
            try:
                # Shape ka main code block (spPr) nikaalo
                spPr = shape._element.find(qn('p:spPr'))
                if spPr is not None:
                    # a) Purane saare colors aur gradients ko delete maaro
                    for fill_type in ['a:solidFill', 'a:gradFill', 'a:pattFill', 'a:blipFill', 'a:noFill']:
                        old_fill = spPr.find(qn(fill_type))
                        if old_fill is not None:
                            spPr.remove(old_fill)
                    
                    # b) Apna naya solid color XML mein inject karo
                    xml_string = f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:srgbClr val="{clean_hex}"/></a:solidFill>'
                    new_fill = parse_xml(xml_string)
                    spPr.append(new_fill)
                    print(f"🔥 XML HACK SUCCESS on: {shape.name}")
            except Exception as xml_e:
                print(f"XML Error on {shape.name}: {xml_e}")

        # 3. Agar Grouped Shape hai, toh uske har bacche (child shape) mein ghus kar rang badlo
        if hasattr(shape, "shapes"):
            for child_shape in shape.shapes:
                apply_theme_color(child_shape, hex_color, force_child=is_target)
                
    except Exception as e:
        pass
# ==========================================
# MAIN GENERATOR FUNCTION
# ==========================================
def generate_pptx(form_data, uploaded_images, uploaded_videos, template_path, output_path):
    prs_source = Presentation(template_path)
    prs_target = Presentation(template_path)
    
    # 1. PEHLE HI SAARE DEFAULT SLIDES DELETE KARDO (Clean Slate)
    for i in range(len(prs_target.slides)-1, -1, -1):
        rId = prs_target.slides._sldIdLst[i].rId
        prs_target.part.drop_rel(rId)
        del prs_target.slides._sldIdLst[i]

    # Theme Color Fetching
    user_theme_color = form_data.get("theme_color", "#0056b3")

    order_raw = form_data.get("slide_order_json", "[]")
    try: chosen_deck = json.loads(order_raw)
    except: chosen_deck = []

    # FIX: 0-Slide Bug Removed. Agar deck khali hai, toh loop nahi chalega.
    
    who_we_are_text = form_data.get("who_we_are", "").strip()
    who_list = [p.strip() for p in who_we_are_text.split('\n') if p.strip()] if who_we_are_text else ["..."]
    our_goals_text = form_data.get("our_goals", "").strip()
    goal_list = [p.strip() for p in our_goals_text.split('\n') if p.strip()] if our_goals_text else ["..."]

    st = {
        "[intro_main]": "HANUAI", "[intro_sub]": "ROADS FOR GROWTH", "[intro_desc]": "Pioneering AI-Driven Solutions",
        "[about_stat1]": form_data.get("about_stat1", ""), "[about_stat2]": form_data.get("about_stat2", ""),
        "[about_stat3]": form_data.get("about_stat3", ""), "[about_stat4]": form_data.get("about_stat4", ""),
        "[erp_main]": form_data.get("erp_main", ""), "[sol_main]": form_data.get("sol_main", ""),
        "[sol_used_main]": form_data.get("sol_used_main", ""), "[pre_post_main]": form_data.get("pre_post_main", ""),
        "[ai_cap_main]": form_data.get("ai_cap_main", ""), "[ai_cap_h1]": form_data.get("ai_cap_h1", ""),
        "[ai_cap_h2]": form_data.get("ai_cap_h2", ""), "[demo_main]": form_data.get("demo_main", ""),
        "[heat_map_main]": form_data.get("heat_map_main", ""), "[cloud_main]": form_data.get("cloud_main", ""),
        "[cloud_sub]": form_data.get("cloud_sub", ""), "[action_main]": form_data.get("action_main", ""),
        "[action_sub]": form_data.get("action_sub", ""), "[work_main]": form_data.get("work_main", ""),
        "[partner_main]": form_data.get("partner_main", ""), "[proj_main]": form_data.get("proj_main", ""),
        "[us_main]": form_data.get("us_main", ""), "[recog_main]": form_data.get("recog_main", ""),
        "[team_main]": form_data.get("team_main", ""), "[thank_main]": form_data.get("thank_main", ""),
        "[thank_sub]": form_data.get("thank_sub", ""), "[thank_copy]": form_data.get("thank_copy", ""),
        "[email_1]": form_data.get("email_1", ""),
        "[Main_heading_intro2]": form_data.get("main_heading_intro2", ""),
        "[Sub_heading_intro2]": form_data.get("sub_heading_intro2", "")
    }

    for i in range(1, 5): 
        st[f"[psh{i}]"] = form_data.get(f"psh_{i}", "")
        st[f"[ds{i}]"] = form_data.get(f"ds_{i}", "")
        st[f"[solh{i}]"] = form_data.get(f"solh_{i}", "")
        st[f"[sold{i}]"] = form_data.get(f"sold_{i}", "")
        st[f"[work_h{i}]"] = form_data.get(f"work_h{i}", "")
        st[f"[work_d{i}]"] = form_data.get(f"work_d{i}", "")
        st[f"[proj_h{i}]"] = form_data.get(f"proj_h{i}", "")
        st[f"[proj_d{i}]"] = form_data.get(f"proj_d{i}", "")
        st[f"[recog_h{i}]"] = form_data.get(f"recog_h_{i}", "") 
        st[f"[recog_d{i}]"] = form_data.get(f"recog_d_{i}", "")
        st[f"[tn{i}]"] = form_data.get(f"team_name_{i}", "")
        st[f"[tr{i}]"] = form_data.get(f"team_role_{i}", "")

    for i in range(1, 9): 
        st[f"[suh{i}]"] = form_data.get(f"suh_{i}", "")
        st[f"[sud{i}]"] = form_data.get(f"sud_{i}", "")
        st[f"[rh{i}]"] = form_data.get(f"recog_h{i}", "")
        st[f"[rd{i}]"] = form_data.get(f"recog_d{i}", "")

    for i in range(1, 8): st[f"[cloud_f{i}]"] = form_data.get(f"cloud_f{i}", "")
    for i in range(1, 6): st[f"[cf{i}]"] = form_data.get(f"cf{i}", "")
    for i in range(1, 7): st[f"[usb{i}]"] = form_data.get(f"us_b{i}", "")

    b_pts = {
        "[ai_cap_b1]": [p.strip() for p in form_data.get("ai_cap_b1", "").split('\n') if p.strip()],
        "[ai_cap_b2]": [p.strip() for p in form_data.get("ai_cap_b2", "").split('\n') if p.strip()],
        "[td1]": [p.strip() for p in form_data.get("team_desc_1", "").split('\n') if p.strip()],
        "[td2]": [p.strip() for p in form_data.get("team_desc_2", "").split('\n') if p.strip()],
        "[td3]": [p.strip() for p in form_data.get("team_desc_3", "").split('\n') if p.strip()],
        "[td4]": [p.strip() for p in form_data.get("team_desc_4", "").split('\n') if p.strip()]
    }

    # Assembly Loop
    for item in chosen_deck:
        slide_id = int(item["slideId"])
        source_idx = slide_id - 1 
        if source_idx < len(prs_source.slides):
            cloned_slide = clone_slide_to_presentation(prs_source, source_idx, prs_target)
            
            for shape in cloned_slide.shapes:
                # 1. THEME COLOR ENGINE 
                apply_theme_color(shape, user_theme_color)

                # 2. TEXT & BULLETS
                txt = get_all_text(shape)
                for tag, val in st.items():
                    if val: replace_text(shape, tag, val)
                if "[who_we_are_1]" in txt: replace_with_bullets(shape, "[who_we_are_1]", who_list)
                elif "[our_goals_1]" in txt: replace_with_bullets(shape, "[our_goals_1]", goal_list)
                for tag, pts in b_pts.items():
                    if tag in txt and pts: replace_with_bullets(shape, tag, pts)

                # 3. IMAGES
                if hasattr(shape, "shape_type") and shape.shape_type == 13:
                    alt = shape._element._nvXxPr.cNvPr.attrib.get('descr', shape.name)
                    img_tags = ["intro_img", "erp_img_1", "erp_img_2", "ai_cap_img_1", "ai_cap_img_2", 
                                "demo_img_1", "heat_map_img_1", "cloud_img", "partner_globe_img",
                                "team_img_1", "team_img_2", "team_img_3", "team_img_4", "thank_img"]
                    for t in range(1, 5): img_tags.append(f"pre_post_img_{t}")
                    for t in img_tags:
                        if t in alt and uploaded_images.get(t):
                            with tempfile.NamedTemporaryFile(delete=False) as tmp:
                                for chunk in uploaded_images[t].chunks(): tmp.write(chunk)
                                tmp_p = tmp.name
                            cloned_slide.shapes.add_picture(tmp_p, shape.left, shape.top, shape.width, shape.height)
                            shape._element.getparent().remove(shape._element)
                            os.unlink(tmp_p)

    if len(prs_target.slides) == 0:
        try: prs_target.slides.add_slide(prs_target.slide_layouts[6])
        except: prs_target.slides.add_slide(prs_target.slide_layouts[0])

    prs_target.save(output_path)